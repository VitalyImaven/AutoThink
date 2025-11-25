"""
File text extraction utilities for various document formats
"""

import json
import xml.etree.ElementTree as ET
from typing import Optional
import io
from app.utils.hebrew_support import normalize_hebrew_text


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        from PyPDF2 import PdfReader
        
        pdf_file = io.BytesIO(file_content)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                # Normalize Hebrew and other Unicode text
                text = normalize_hebrew_text(text)
                text_parts.append(text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file"""
    try:
        from docx import Document
        
        doc_file = io.BytesIO(file_content)
        doc = Document(doc_file)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from XLSX file"""
    try:
        from openpyxl import load_workbook
        
        excel_file = io.BytesIO(file_content)
        workbook = load_workbook(excel_file, read_only=True, data_only=True)
        
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from XLSX: {str(e)}")


def extract_text_from_json(file_content: bytes) -> str:
    """Extract text from JSON file"""
    try:
        content = file_content.decode('utf-8')
        data = json.loads(content)
        
        # Pretty print JSON as text
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        raise ValueError(f"Failed to extract text from JSON: {str(e)}")


def extract_text_from_xml(file_content: bytes) -> str:
    """Extract text from XML file"""
    try:
        content = file_content.decode('utf-8')
        root = ET.fromstring(content)
        
        # Extract all text nodes
        text_parts = []
        
        def extract_text(element, prefix=""):
            if element.text and element.text.strip():
                text_parts.append(f"{prefix}{element.tag}: {element.text.strip()}")
            
            for child in element:
                extract_text(child, prefix + "  ")
        
        extract_text(root)
        return "\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from XML: {str(e)}")


def extract_text_from_markdown(file_content: bytes) -> str:
    """Extract text from Markdown file"""
    try:
        # Markdown is already text, just decode
        return file_content.decode('utf-8')
    except Exception as e:
        raise ValueError(f"Failed to extract text from Markdown: {str(e)}")


def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from TXT file"""
    try:
        # Try UTF-8 first, fallback to latin-1
        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            return file_content.decode('latin-1')
    except Exception as e:
        raise ValueError(f"Failed to extract text from TXT: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file"""
    try:
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        text_parts = []
        for i, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {str(e)}")


def extract_text_from_file(filename: str, file_content: bytes) -> str:
    """
    Extract text from various file formats based on file extension
    
    Args:
        filename: Name of the file
        file_content: Binary content of the file
        
    Returns:
        Extracted text content
        
    Raises:
        ValueError: If file type is not supported or extraction fails
    """
    extension = filename.lower().split('.')[-1]
    
    extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'doc': extract_text_from_docx,  # Try docx extractor for .doc
        'xlsx': extract_text_from_xlsx,
        'xls': extract_text_from_xlsx,  # Try xlsx extractor for .xls
        'json': extract_text_from_json,
        'xml': extract_text_from_xml,
        'md': extract_text_from_markdown,
        'markdown': extract_text_from_markdown,
        'txt': extract_text_from_txt,
        'text': extract_text_from_txt,
        'pptx': extract_text_from_pptx,
        'ppt': extract_text_from_pptx,
    }
    
    extractor = extractors.get(extension)
    if not extractor:
        raise ValueError(
            f"Unsupported file type: .{extension}. "
            f"Supported types: {', '.join(sorted(extractors.keys()))}"
        )
    
    return extractor(file_content)

